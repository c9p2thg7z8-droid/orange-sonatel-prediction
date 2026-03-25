from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

ORANGE = RGBColor(0xFF, 0x66, 0x00)
TEAL   = RGBColor(0x00, 0x80, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)

LOGO   = "logo.png"
SC_BIE = "screenshots/bienvenue.png"
SC_INT = "screenshots/interface.png"
SC_RES = "screenshots/resultat.png"

blank = prs.slide_layouts[6]  # layout vide

def add_rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_logo(slide, x, y, h=0.7):
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(x), Inches(y), height=Inches(h))

# ─────────────────────────────────────────
# SLIDE 1 — TITRE
# ─────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, ORANGE)
add_rect(sl, 0, 5.2, 13.33, 2.3, TEAL)
# Cercles décoratifs
for cx,cy,cw,ch in [(10,-.5,4,4),(-.5,4.5,3,3)]:
    s = sl.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(cw), Inches(ch))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0xFF,0x88,0x00)
    s.line.fill.background()

# Logo
if os.path.exists(LOGO):
    sl.shapes.add_picture(LOGO, Inches(0.5), Inches(0.3), height=Inches(1.0))

add_text(sl, "Système de Prédiction des Groupes", 0.5, 1.8, 12, 1.2, 36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Affectation automatique des plaintes clients", 0.5, 3.0, 12, 0.7, 20, color=RGBColor(0xFF,0xEE,0xDD), align=PP_ALIGN.CENTER)
add_rect(sl, 4.5, 3.9, 4.3, 0.9, WHITE)
add_text(sl, "DSI — Direction des Systèmes d'Information", 4.6, 4.0, 4.1, 0.7, 13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(sl, "Orange Sonatel  •  2025", 0.5, 6.8, 12, 0.5, 12, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────────────────────────
# SLIDE 2 — CONTEXTE
# ─────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
add_rect(sl, 0, 0, 13.33, 1.1, ORANGE)
add_logo(sl, 0.3, 0.2)
add_text(sl, "Contexte & Problématique", 1.5, 0.15, 11, 0.8, 22, bold=True, color=WHITE)

items = [
    ("📌", "Contexte", "Stage à la DSI d'Orange Sonatel — traitement des plaintes clients"),
    ("⚠️", "Problème", "Affectation manuelle des plaintes aux groupes de traitement : lente et sujette aux erreurs"),
    ("🎯", "Objectif", "Automatiser l'affectation grâce au Machine Learning pour gagner en efficacité"),
    ("💡", "Solution", "Modèle Random Forest entraîné sur l'historique des plaintes classifiées"),
]
for i,(icon,titre,desc) in enumerate(items):
    y = 1.3 + i*1.4
    add_rect(sl, 0.4, y, 12.5, 1.2, WHITE)
    add_rect(sl, 0.4, y, 0.5, 1.2, ORANGE)
    add_text(sl, icon, 0.45, y+0.25, 0.5, 0.7, 18, align=PP_ALIGN.CENTER)
    add_text(sl, titre, 1.1, y+0.05, 3, 0.45, 13, bold=True, color=ORANGE)
    add_text(sl, desc,  1.1, y+0.5,  11, 0.55, 12, color=DARK)

# ─────────────────────────────────────────
# SLIDE 3 — DONNÉES & MODÈLE
# ─────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
add_rect(sl, 0, 0, 13.33, 1.1, TEAL)
add_logo(sl, 0.3, 0.2)
add_text(sl, "Données & Modèle ML", 1.5, 0.15, 11, 0.8, 22, bold=True, color=WHITE)

# Colonne gauche
add_rect(sl, 0.4, 1.3, 5.8, 5.7, WHITE)
add_rect(sl, 0.4, 1.3, 5.8, 0.5, TEAL)
add_text(sl, "📊  Variables d'entrée", 0.6, 1.35, 5.4, 0.4, 13, bold=True, color=WHITE)
vars_ = ["DOMAINE","SOUS-DOMAINE","TYPE","TYPOLOGIE"]
for i,v in enumerate(vars_):
    add_rect(sl, 0.6, 2.0+i*1.1, 5.4, 0.85, LGRAY)
    add_rect(sl, 0.6, 2.0+i*1.1, 0.08, 0.85, TEAL)
    add_text(sl, v, 0.85, 2.1+i*1.1, 5, 0.6, 13, bold=True, color=DARK)

# Colonne droite
add_rect(sl, 6.8, 1.3, 6.1, 5.7, WHITE)
add_rect(sl, 6.8, 1.3, 6.1, 0.5, ORANGE)
add_text(sl, "🤖  Modèle & Performance", 7.0, 1.35, 5.7, 0.4, 13, bold=True, color=WHITE)
infos = [
    ("Algorithme","Random Forest Classifier"),
    ("Librairie","scikit-learn 1.6.1"),
    ("Encodage","LabelEncoder par variable"),
    ("Sortie","Groupe prédit + % confiance"),
]
for i,(k,v) in enumerate(infos):
    y = 2.0+i*1.1
    add_rect(sl, 7.0, y, 5.7, 0.85, LGRAY)
    add_text(sl, k, 7.1, y+0.05, 2.2, 0.4, 11, bold=True, color=ORANGE)
    add_text(sl, v, 7.1, y+0.42, 5.4, 0.38, 11, color=DARK)

# ─────────────────────────────────────────
# SLIDE 4 — ÉCRAN DE BIENVENUE
# ─────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
add_rect(sl, 0, 0, 13.33, 1.1, ORANGE)
add_logo(sl, 0.3, 0.2)
add_text(sl, "Interface — Écran de Bienvenue", 1.5, 0.15, 11, 0.8, 22, bold=True, color=WHITE)
if os.path.exists(SC_BIE):
    sl.shapes.add_picture(SC_BIE, Inches(1.5), Inches(1.3), width=Inches(10.3))
add_text(sl, "Splash screen animé avec logo Orange Sonatel, titre DSI et points de chargement", 0.5, 6.9, 12.3, 0.5, 11, color=DARK, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────────────────────────
# SLIDE 5 — INTERFACE PRINCIPALE
# ─────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
add_rect(sl, 0, 0, 13.33, 1.1, TEAL)
add_logo(sl, 0.3, 0.2)
add_text(sl, "Interface — Prédiction Simple", 1.5, 0.15, 11, 0.8, 22, bold=True, color=WHITE)
if os.path.exists(SC_INT):
    sl.shapes.add_picture(SC_INT, Inches(1.5), Inches(1.3), width=Inches(10.3))
add_text(sl, "Sélection des paramètres via dropdowns avec autocomplétion — valeurs issues du modèle", 0.5, 6.9, 12.3, 0.5, 11, color=DARK, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────────────────────────
# SLIDE 6 — RÉSULTAT
# ─────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
add_rect(sl, 0, 0, 13.33, 1.1, ORANGE)
add_logo(sl, 0.3, 0.2)
add_text(sl, "Interface — Résultat de Prédiction", 1.5, 0.15, 11, 0.8, 22, bold=True, color=WHITE)
if os.path.exists(SC_RES):
    sl.shapes.add_picture(SC_RES, Inches(1.5), Inches(1.3), width=Inches(10.3))
add_text(sl, "Carte animée avec groupe prédit, barre de progression et niveau de confiance coloré", 0.5, 6.9, 12.3, 0.5, 11, color=DARK, align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────────────────────────
# SLIDE 7 — FONCTIONNALITÉS
# ─────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
add_rect(sl, 0, 0, 13.33, 1.1, TEAL)
add_logo(sl, 0.3, 0.2)
add_text(sl, "Fonctionnalités", 1.5, 0.15, 11, 0.8, 22, bold=True, color=WHITE)

feats = [
    ("🎯","Prédiction simple","Dropdowns avec autocomplétion, résultat instantané avec niveau de confiance"),
    ("📂","Prédiction en masse","Upload CSV, traitement batch, téléchargement des résultats"),
    ("📊","Résultat visuel","Carte animée, barre de progression colorée selon la confiance"),
    ("ℹ️","Infos modèle","Détails techniques : variables, groupes disponibles, statistiques"),
    ("🌐","Site web","Interface HTML/CSS/JS complète via API FastAPI"),
    ("🔒","Validation","Détection des valeurs inconnues avec message d'avertissement"),
]
for i,(icon,titre,desc) in enumerate(feats):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.4 + row * 1.9
    add_rect(sl, x, y, 6.1, 1.6, WHITE)
    add_rect(sl, x, y, 0.6, 1.6, ORANGE if col==0 else TEAL)
    add_text(sl, icon,  x+0.1, y+0.4, 0.5, 0.8, 18, align=PP_ALIGN.CENTER)
    add_text(sl, titre, x+0.75, y+0.1, 5.2, 0.5, 13, bold=True, color=ORANGE if col==0 else TEAL)
    add_text(sl, desc,  x+0.75, y+0.65, 5.2, 0.8, 11, color=DARK)

# ─────────────────────────────────────────
# SLIDE 8 — STACK TECHNIQUE
# ─────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, LGRAY)
add_rect(sl, 0, 0, 13.33, 1.1, ORANGE)
add_logo(sl, 0.3, 0.2)
add_text(sl, "Stack Technique", 1.5, 0.15, 11, 0.8, 22, bold=True, color=WHITE)

techs = [
    ("🐍","Python","Langage principal"),
    ("🤖","scikit-learn","Modèle Random Forest"),
    ("🎨","Gradio","Interface web"),
    ("🐼","pandas","Traitement données"),
    ("⚡","FastAPI","API REST"),
    ("🌐","HTML/CSS/JS","Site web frontend"),
]
for i,(icon,nom,usage) in enumerate(techs):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.2
    y = 1.4 + row * 2.6
    add_rect(sl, x, y, 3.8, 2.2, WHITE)
    add_rect(sl, x, y, 3.8, 0.55, ORANGE if row==0 else TEAL)
    add_text(sl, icon+" "+nom, x+0.15, y+0.07, 3.5, 0.42, 14, bold=True, color=WHITE)
    add_text(sl, usage, x+0.15, y+0.75, 3.5, 1.2, 12, color=DARK)

# ─────────────────────────────────────────
# SLIDE 9 — CONCLUSION
# ─────────────────────────────────────────
sl = prs.slides.add_slide(blank)
add_rect(sl, 0, 0, 13.33, 7.5, ORANGE)
add_rect(sl, 0, 5.5, 13.33, 2.0, TEAL)
for cx,cy,cw,ch in [(9.5,-.5,5,5),(-.5,5,4,4)]:
    s = sl.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(cw), Inches(ch))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0xFF,0x88,0x00)
    s.line.fill.background()
if os.path.exists(LOGO):
    sl.shapes.add_picture(LOGO, Inches(0.5), Inches(0.3), height=Inches(0.9))
add_text(sl, "Merci pour votre attention", 0.5, 1.8, 12.3, 1.0, 34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Système de Prédiction des Groupes — Orange Sonatel", 0.5, 3.0, 12.3, 0.7, 16, color=RGBColor(0xFF,0xEE,0xDD), align=PP_ALIGN.CENTER)
add_rect(sl, 4.0, 3.9, 5.3, 1.0, WHITE)
add_text(sl, "Cheikh FAYE — Stagiaire DSI", 4.1, 4.0, 5.1, 0.45, 13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(sl, "© 2025 Orange Sonatel", 4.1, 4.45, 5.1, 0.4, 11, color=TEAL, align=PP_ALIGN.CENTER)
add_text(sl, "github.com/c9p2thg7z8-droid/orange-sonatel-prediction", 0.5, 6.8, 12.3, 0.5, 11, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

prs.save("presentation_orange_sonatel.pptx")
print("✅ Fichier généré : presentation_orange_sonatel.pptx")
